from __future__ import annotations

import csv
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt


BASE_DIR = Path(__file__).resolve().parent
TABLES_DIR = BASE_DIR / "Tables"
FIGURES_DIR = BASE_DIR / "Figures"
OUTPUT_PATH = BASE_DIR / "Heart_Failure_Risk_Presentation.pptx"


def read_csv(path: Path) -> list[list[str]]:
    with path.open(newline="", encoding="utf-8") as file:
        return list(csv.reader(file))


def add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle


def add_bullets_slide(prs: Presentation, title: str, bullets: list[str], subtitle: str | None = None) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    text_frame = slide.placeholders[1].text_frame
    text_frame.clear()
    if subtitle:
        para = text_frame.paragraphs[0]
        para.text = subtitle
        para.font.bold = True
        para.font.size = Pt(20)
    else:
        text_frame.paragraphs[0].text = bullets[0]
        bullets = bullets[1:]

    for bullet in bullets:
        para = text_frame.add_paragraph()
        para.text = bullet
        para.level = 0
        para.font.size = Pt(20)


def add_table_slide(prs: Presentation, title: str, rows: list[list[str]], top_n: int | None = None) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title

    data = rows if top_n is None else [rows[0], *rows[1 : top_n + 1]]
    n_rows, n_cols = len(data), len(data[0])

    table_shape = slide.shapes.add_table(n_rows, n_cols, Inches(0.4), Inches(1.4), Inches(12.5), Inches(5.3))
    table = table_shape.table

    for r in range(n_rows):
        for c in range(n_cols):
            cell = table.cell(r, c)
            cell.text = data[r][c]
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(12 if r == 0 else 11)
                if r == 0:
                    paragraph.font.bold = True


def add_image_slide(prs: Presentation, title: str, image_path: Path, caption: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    slide.shapes.add_picture(str(image_path), Inches(0.8), Inches(1.4), width=Inches(11.7))
    caption_box = slide.shapes.add_textbox(Inches(0.8), Inches(6.6), Inches(11.7), Inches(0.6))
    caption_box.text_frame.text = caption
    caption_box.text_frame.paragraphs[0].font.size = Pt(16)


def add_agenda_slide(prs: Presentation) -> None:
    add_bullets_slide(
        prs,
        "Agenda",
        [
            "Clinical motivation and problem statement",
            "Dataset and preprocessing",
            "Model development and hybrid risk score",
            "Performance and threshold optimization",
            "Explainability, counterfactuals, and conclusions",
        ],
    )


def create_presentation() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    add_title_slide(
        prs,
        "Heart Failure Risk Score Prediction",
        "Major Project Presentation\nPrepared for Professor Review",
    )
    add_agenda_slide(prs)

    dataset_rows = read_csv(TABLES_DIR / "T1_dataset_overview.csv")
    training_rows = read_csv(TABLES_DIR / "model_training_summary.csv")
    perf_rows = read_csv(TABLES_DIR / "Table11__Model_Performance_Comparison.csv")
    threshold_rows = read_csv(TABLES_DIR / "Table14_Threshold_Optimization.csv")
    summary_rows = read_csv(TABLES_DIR / "Table15_Final_Summary.csv")
    counterfactual_rows = read_csv(TABLES_DIR / "Counterfactual_Summary_updated.csv")

    add_bullets_slide(
        prs,
        "Problem Statement & Objective",
        [
            "Heart disease remains a leading cause of mortality worldwide.",
            "Goal: build robust, interpretable models to predict heart failure risk.",
            "Evaluate ML models and create a Hybrid Risk Score (HRS) for clinical use.",
            "Optimize decision thresholds for stronger sensitivity/specificity balance.",
        ],
    )

    add_table_slide(prs, "Dataset Overview", dataset_rows)

    add_image_slide(
        prs,
        "Data Exploration: Class Distribution",
        FIGURES_DIR / "class_distribution.png",
        "The dataset has class imbalance that requires careful evaluation beyond accuracy.",
    )

    add_image_slide(
        prs,
        "Feature Relationships",
        FIGURES_DIR / "correlation_heatmap.png",
        "Correlation analysis guided model and feature interpretation steps.",
    )

    add_table_slide(prs, "Model Training Configuration", training_rows)

    add_table_slide(prs, "Model Performance Comparison", perf_rows)

    add_image_slide(
        prs,
        "ROC Curve Comparison",
        FIGURES_DIR / "Plot12_ROC_Curves (1).png",
        "CatBoost and XGBoost achieved the highest discriminatory power.",
    )

    add_image_slide(
        prs,
        "Confusion Matrix Comparison",
        FIGURES_DIR / "Plot12b_Confusion_Matrices.png",
        "Confusion matrices show class-wise trade-offs across models.",
    )

    add_image_slide(
        prs,
        "Hybrid Risk Score Distribution",
        FIGURES_DIR / "hrs_score_distribution.png",
        "HRS provides a clinically interpretable stratification framework.",
    )

    add_table_slide(prs, "Threshold Optimization (Aurobindo's Contribution)", threshold_rows)

    add_image_slide(
        prs,
        "Calibration Curve",
        FIGURES_DIR / "Plot15_Calibration_Curve.png",
        "Calibration quality shows how closely predicted risk matches observed outcomes.",
    )

    add_image_slide(
        prs,
        "Decision Curve Analysis",
        FIGURES_DIR / "Plot16_Decision_Curve_Analysis.png",
        "DCA indicates net clinical benefit across probability thresholds.",
    )

    add_image_slide(
        prs,
        "Threshold Optimization Curves",
        FIGURES_DIR / "Plot18_Threshold_Optimization_Curves.png",
        "Optimal thresholds selected using Youden J-statistics and clinical balance.",
    )

    add_table_slide(prs, "Final Model Summary", summary_rows)

    add_image_slide(
        prs,
        "SHAP Global Interpretability (XGBoost)",
        FIGURES_DIR / "xgboost_shap_beeswarm.png",
        "SHAP highlights global feature effects and model transparency.",
    )

    add_image_slide(
        prs,
        "Counterfactual Impact",
        FIGURES_DIR / "Counterfactual_Impact_updated.png",
        "Counterfactuals suggest actionable feature changes to reduce individual risk.",
    )

    add_table_slide(prs, "Counterfactual Summary (Sample)", counterfactual_rows, top_n=8)

    add_bullets_slide(
        prs,
        "Team Contributions",
        [
            "Arpit: data preprocessing and quality checks.",
            "Aditya: model evaluation (accuracy, precision, recall, F1, AUC).",
            "Aurobindo: threshold optimization, calibration, Brier score, DCA, and final integration.",
            "Collaborative outcome: explainable and clinically oriented risk prediction workflow.",
        ],
    )

    add_bullets_slide(
        prs,
        "Conclusion & Future Work",
        [
            "CatBoost and XGBoost delivered strongest predictive performance.",
            "Threshold tuning improved real-world decision utility.",
            "HRS and explainability improve interpretability for clinicians.",
            "Future work: external validation, prospective testing, and deployment as a decision-support tool.",
        ],
    )

    add_bullets_slide(
        prs,
        "Q&A",
        [
            "Thank you.",
            "I would be happy to discuss methodology, assumptions, and clinical implications.",
        ],
    )

    prs.save(OUTPUT_PATH)
    print(f"Saved: {OUTPUT_PATH}")


if __name__ == "__main__":
    create_presentation()
